import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Premium Dark Mode Theme Colors
    BG_DARK = RGBColor(10, 25, 47)         # Deep space navy background
    CARD_BG = RGBColor(23, 42, 69)         # Lighter navy for cards
    WHITE = RGBColor(255, 255, 255)        # Pure white for main titles
    SILVER = RGBColor(203, 213, 225)       # Muted silver for body text
    CYAN = RGBColor(0, 240, 255)           # Electric cyan for primary accents
    PURPLE = RGBColor(139, 92, 246)        # Royal purple for secondary accents
    RED_ACCENT = RGBColor(244, 63, 94)     # Red for CFT vulnerability
    GREEN_ACCENT = RGBColor(52, 211, 153)  # Green for BFT resilience
    DARK_GRAY = RGBColor(148, 163, 184)    # Light slate gray (improved readability)

    blank_layout = prs.slide_layouts[6]

    def add_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def draw_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CYAN, border_width=1.5):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
        return shape

    def add_header(slide, title, category="PROJET DE FIN D'ANNÉE — EMSI"):
        # Category tracker
        catBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.25), Inches(11.83), Inches(0.4))
        ctf = catBox.text_frame
        ctf.margin_top = ctf.margin_bottom = ctf.margin_left = ctf.margin_right = 0
        cp = ctf.paragraphs[0]
        cp.text = category.upper()
        cp.font.name = 'Inter'
        cp.font.size = Pt(11)
        cp.font.bold = True
        cp.font.color.rgb = PURPLE
        
        # Main Title
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.83), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Inter'
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Glowing accent line below title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(11.83), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = CYAN
        line.line.fill.background()

    def set_text_box_formatting(tf):
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

    # ============================================================
    # SLIDE 1 : Page de garde (Sleek EMSI Dark Theme)
    # ============================================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1)
    
    # Large futuristic layout back card
    draw_card(slide1, Inches(0.75), Inches(0.5), Inches(11.83), Inches(6.5), bg_color=CARD_BG, border_color=PURPLE, border_width=2)
    
    # School header
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "ÉCOLE MAROCAINE DES SCIENCES DE L'INGÉNIEUR"
    p1.font.name = 'Inter'
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(6)
    
    p2 = tf.add_paragraph()
    p2.text = "Membre de Honoris United Universities"
    p2.font.name = 'Inter'
    p2.font.size = Pt(14)
    p2.font.italic = True
    p2.font.color.rgb = CYAN
    p2.alignment = PP_ALIGN.CENTER
    
    # Title
    txBox2 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = "PROJET DE FIN D'ANNÉE (PFA)"
    p3.font.name = 'Inter'
    p3.font.size = Pt(24)
    p3.font.bold = True
    p3.font.color.rgb = SILVER
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(10)
    
    p4 = tf2.add_paragraph()
    p4.text = "Implémentation d'un Système de Consensus Byzantin Résilient (PBFT)\net Sécurisation Cryptographique"
    p4.font.name = 'Inter'
    p4.font.size = Pt(32)
    p4.font.bold = True
    p4.font.color.rgb = CYAN
    p4.alignment = PP_ALIGN.CENTER
    p4.line_spacing = 1.2
    
    # Glowing separator
    sep = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(4.5), Inches(5.33), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = PURPLE
    sep.line.fill.background()

    # Author details Left Box
    txBox3 = slide1.shapes.add_textbox(Inches(1.5), Inches(4.9), Inches(4.5), Inches(1.8))
    tf3 = txBox3.text_frame
    p5 = tf3.paragraphs[0]
    p5.text = "Réalisé par :"
    p5.font.name = 'Inter'
    p5.font.size = Pt(15)
    p5.font.bold = True
    p5.font.color.rgb = CYAN
    p5.space_after = Pt(8)
    
    for author in ["Yassir KADOUARI", "Matine ELKASBIJI", "Marouane ISMAILI"]:
        p_auth = tf3.add_paragraph()
        p_auth.text = f"• {author}"
        p_auth.font.name = 'Inter'
        p_auth.font.size = Pt(13)
        p_auth.font.color.rgb = WHITE
        p_auth.space_after = Pt(4)
        
    # Info details Right Box
    txBox4 = slide1.shapes.add_textbox(Inches(7.5), Inches(4.9), Inches(4.5), Inches(1.8))
    tf4 = txBox4.text_frame
    p6 = tf4.paragraphs[0]
    p6.text = "Filière & Option :"
    p6.font.name = 'Inter'
    p6.font.size = Pt(15)
    p6.font.bold = True
    p6.font.color.rgb = PURPLE
    p6.space_after = Pt(8)
    
    p_fil = tf4.add_paragraph()
    p_fil.text = "Ingénierie Informatique\nSystèmes, Réseaux & Génie Logiciel"
    p_fil.font.name = 'Inter'
    p_fil.font.size = Pt(13)
    p_fil.font.color.rgb = WHITE
    p_fil.line_spacing = 1.15
    p_fil.space_after = Pt(8)
    
    p_year = tf4.add_paragraph()
    p_year.text = "Année Universitaire : 2025 / 2026"
    p_year.font.name = 'Inter'
    p_year.font.size = Pt(12)
    p_year.font.color.rgb = SILVER

    # ============================================================
    # SLIDE 2 : Contexte & Enjeux (Futuristic Cards Layout - 10+ lines)
    # ============================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2)
    add_header(slide2, "1. Contexte & Enjeux Industriels")
    
    # Left Card: Distributed Systems Challenges
    draw_card(slide2, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3), border_color=CYAN)
    txBox_l = slide2.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_l = txBox_l.text_frame
    set_text_box_formatting(tf_l)
    p_l = tf_l.paragraphs[0]
    p_l.text = "Enjeux des Systèmes Distribués"
    p_l.font.name = 'Inter'
    p_l.font.size = Pt(18)
    p_l.font.bold = True
    p_l.font.color.rgb = CYAN
    p_l.space_after = Pt(14)
    
    pts_l = [
        "Infrastructures hautement critiques globales.",
        "Transactions bancaires, chaînes de blocs, clouds.",
        "Besoin de coordination sur un état de registre partagé.",
        "Pannes matérielles franches (coupure de courant, crash).",
        "Saturations réseau asynchrones et latences imprévues."
    ]
    for pt in pts_l:
        p_pt = tf_l.add_paragraph()
        p_pt.text = f"→ {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = WHITE
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # Right Card: The Byzantine Threat
    draw_card(slide2, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3), border_color=PURPLE)
    txBox_r = slide2.shapes.add_textbox(Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_r = txBox_r.text_frame
    set_text_box_formatting(tf_r)
    p_r = tf_r.paragraphs[0]
    p_r.text = "Le Modèle de Panne Byzantine"
    p_r.font.name = 'Inter'
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = PURPLE
    p_r.space_after = Pt(14)
    
    pts_r = [
        "Introduction de nœuds piratés ou de bugs arbitraires.",
        "Omissions volontaires, falsifications actives de paquets.",
        "Équivocations : diffusion de messages contradictoires.",
        "Formule de tolérance robuste : N >= 3f + 1.",
        "Notre cluster à 4 nœuds tolère précisément f = 1 panne."
    ]
    for pt in pts_r:
        p_pt = tf_r.add_paragraph()
        p_pt.text = f"→ {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = SILVER
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # ============================================================
    # SLIDE 3 : CFT vs BFT (Side-by-Side Glowing Panels - 10+ lines)
    # ============================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3)
    add_header(slide3, "2. Problématique : CFT vs BFT")
    
    # Left Card: CFT
    draw_card(slide3, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3), border_color=RED_ACCENT)
    txBox_l = slide3.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_l = txBox_l.text_frame
    set_text_box_formatting(tf_l)
    p_l = tf_l.paragraphs[0]
    p_l.text = "Crash Fault Tolerant (CFT)"
    p_l.font.name = 'Inter'
    p_l.font.size = Pt(18)
    p_l.font.bold = True
    p_l.font.color.rgb = RED_ACCENT
    p_l.space_after = Pt(14)
    
    cft_points = [
        "Modèle restreint : crash franc et arrêts nets.",
        "Solutions industrielles : Raft et Paxos standard.",
        "Hypothèse interne : aucun nœud ne triche ou ne ment.",
        "Tolérance : uniquement coupures temporaires de courant.",
        "❌ Échouent dès qu'un nœud est piraté ou malicieux.",
        "Usage : Bases de données internes de confiance."
    ]
    for pt in cft_points:
        p_pt = tf_l.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = SILVER
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # Right Card: BFT
    draw_card(slide3, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3), border_color=GREEN_ACCENT)
    txBox_r = slide3.shapes.add_textbox(Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_r = txBox_r.text_frame
    set_text_box_formatting(tf_r)
    p_r = tf_r.paragraphs[0]
    p_r.text = "Byzantine Fault Tolerant (BFT)"
    p_r.font.name = 'Inter'
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = GREEN_ACCENT
    p_r.space_after = Pt(14)
    
    bft_points = [
        "Modèle total : pannes actives, logiques et arbitraires.",
        "Algorithmes phares : PBFT, HotStuff, Tendermint.",
        "Hypothèse interne : certains réplicas mentent activement.",
        "Tolérance : falsifications de votes, replay attacks.",
        "✅ Maintient l'état sain via une double phase de quorum.",
        "Usage : Blockchains ouvertes, coordination inter-bancaire."
    ]
    for pt in bft_points:
        p_pt = tf_r.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = WHITE
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # ============================================================
    # SLIDE 4 : Le Protocole PBFT (Timeline of 3 Cards - 10+ lines)
    # ============================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4)
    add_header(slide4, "3. Fonctionnement du Consensus PBFT")
    
    # 3 Cards Side-by-Side representing the 3 phases
    card_w = Inches(3.7)
    card_h = Inches(5.3)
    
    # Card 1: Pre-Prepare
    draw_card(slide4, Inches(0.75), Inches(1.6), card_w, card_h, border_color=CYAN)
    tx_1 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.6), card_w, card_h)
    tf_1 = tx_1.text_frame
    set_text_box_formatting(tf_1)
    p_1 = tf_1.paragraphs[0]
    p_1.text = "Phase 1 : Pre-Prepare"
    p_1.font.name = 'Inter'
    p_1.font.size = Pt(16)
    p_1.font.bold = True
    p_1.font.color.rgb = CYAN
    p_1.space_after = Pt(14)
    
    pts_1 = [
        "Réception : Le Leader reçoit la transaction cliente.",
        "Séquence : Le leader affecte un numéro séquentiel strict.",
        "Signature : Calcul et injection d'une signature ECDSA.",
        "Diffusion : Envoi de la proposition à tout le cluster."
    ]
    for pt in pts_1:
        p_pt = tf_1.add_paragraph()
        p_pt.text = f"→ {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = WHITE
        p_pt.space_after = Pt(12)
        p_pt.line_spacing = 1.15

    # Card 2: Prepare
    draw_card(slide4, Inches(4.81), Inches(1.6), card_w, card_h, border_color=PURPLE)
    tx_2 = slide4.shapes.add_textbox(Inches(4.81), Inches(1.6), card_w, card_h)
    tf_2 = tx_2.text_frame
    set_text_box_formatting(tf_2)
    p_2 = tf_2.paragraphs[0]
    p_2.text = "Phase 2 : Prepare"
    p_2.font.name = 'Inter'
    p_2.font.size = Pt(16)
    p_2.font.bold = True
    p_2.font.color.rgb = PURPLE
    p_2.space_after = Pt(14)
    
    pts_2 = [
        "Vérification : Les réplicas vérifient la signature du leader.",
        "Vote : Chaque nœud sain émet un vote PREPARE.",
        "Quorum : Accumulation de 2f votes valides distincts.",
        "État Prepared : Transition locale à l'état stable Prepared."
    ]
    for pt in pts_2:
        p_pt = tf_2.add_paragraph()
        p_pt.text = f"→ {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = SILVER
        p_pt.space_after = Pt(12)
        p_pt.line_spacing = 1.15

    # Card 3: Commit
    draw_card(slide4, Inches(8.87), Inches(1.6), card_w, card_h, border_color=GREEN_ACCENT)
    tx_3 = slide4.shapes.add_textbox(Inches(8.87), Inches(1.6), card_w, card_h)
    tf_3 = tx_3.text_frame
    set_text_box_formatting(tf_3)
    p_3 = tf_3.paragraphs[0]
    p_3.text = "Phase 3 : Commit"
    p_3.font.name = 'Inter'
    p_3.font.size = Pt(16)
    p_3.font.bold = True
    p_3.font.color.rgb = GREEN_ACCENT
    p_3.space_after = Pt(14)
    
    pts_3 = [
        "Diffusion : Les nœuds prepared s'envoient un vote COMMIT.",
        "Double Quorum : Collecte de 2f + 1 commits valides.",
        "Exécution : Application définitive dans le registre local.",
        "Réponse client : Retour d'un message REPLY signé."
    ]
    for pt in pts_3:
        p_pt = tf_3.add_paragraph()
        p_pt.text = f"→ {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = WHITE
        p_pt.space_after = Pt(12)
        p_pt.line_spacing = 1.15

    # ============================================================
    # SLIDE 5 : Cryptographie ECDSA & Shamir (Stunning Panels - 10+ lines)
    # ============================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5)
    add_header(slide5, "4. Cryptographie Asymétrique & Seuil")
    
    # Left Card: ECDSA
    draw_card(slide5, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3), border_color=CYAN)
    txBox_l = slide5.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_l = txBox_l.text_frame
    set_text_box_formatting(tf_l)
    p_l = tf_l.paragraphs[0]
    p_l.text = "Signatures ECDSA (secp256r1)"
    p_l.font.name = 'Inter'
    p_l.font.size = Pt(18)
    p_l.font.bold = True
    p_l.font.color.rgb = CYAN
    p_l.space_after = Pt(14)
    
    ecdsa_pts = [
        "Briques d'authentification basées sur BouncyCastle.",
        "Handshake initial : échange automatisé des clés publiques.",
        "Garantie : signatures uniques jointes aux paquets Prepare/Commit.",
        "Vérification instantanée en arithmétique modulaire.",
        "Sécurité absolue : aucune falsification possible des votes."
    ]
    for pt in ecdsa_pts:
        p_pt = tf_l.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = WHITE
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # Right Card: Shamir
    draw_card(slide5, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3), border_color=PURPLE)
    txBox_r = slide5.shapes.add_textbox(Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_r = txBox_r.text_frame
    set_text_box_formatting(tf_r)
    p_r = tf_r.paragraphs[0]
    p_r.text = "Partage de Secret de Shamir (t=2, n=4)"
    p_r.font.name = 'Inter'
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = PURPLE
    p_r.space_after = Pt(14)
    
    shamir_pts = [
        "Partage de secret cryptographique de groupe décentralisé.",
        "Clé répartie géométriquement en n = 4 parts (shares) distinctes.",
        "Reconstruction possible dès que t = 2 nœuds sains collaborent.",
        "Interpolation de Lagrange modulo le grand nombre premier P.",
        "Architecture décentralisée éliminant les points de faille."
    ]
    for pt in shamir_pts:
        p_pt = tf_r.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = SILVER
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # ============================================================
    # SLIDE 6 : Architecture Logicielle (2x2 Grid Layout spaced out beautifully)
    # ============================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6)
    add_header(slide6, "5. Architecture & Développement Système")
    
    cw = Inches(5.6)
    ch = Inches(2.4) # Slightly shorter cards to avoid overflowing with margins
    
    grid = [
        (Inches(0.75), Inches(1.6), "NetworkManager (TCP Mesh)", "Sockets asynchrones multi-threadées / maillage complet / protection anti-relecture SHA-256.", CYAN),
        (Inches(6.98), Inches(1.6), "PBFTEngine (Automate PBFT)", "Gestion des quorums et des états Prepared/Committed / machine d'état asynchrone.", PURPLE),
        (Inches(0.75), Inches(4.3), "ByzantineBehavior (Anomalies)", "Injection à chaud d'attaques logiques (SilentNode, Equivocation, Replay).", GREEN_ACCENT),
        (Inches(6.98), Inches(4.3), "MetricsServer (HTTP REST API)", "Serveur Sun HttpServer intégré / exposition instantanée en JSON.", WHITE)
    ]
    
    for left, top, title, desc, color in grid:
        draw_card(slide6, left, top, cw, ch, border_color=color, border_width=1.2)
        tx = slide6.shapes.add_textbox(left, top, cw, ch)
        tf = tx.text_frame
        set_text_box_formatting(tf)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Inter'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(8) # Nice spacing between card title and bullets
        
        p_d = tf.add_paragraph()
        p_d.text = f"→ {desc}"
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = SILVER
        p_d.line_spacing = 1.15
        p_d.space_after = Pt(6)
        
        p_d2 = tf.add_paragraph()
        p_d2.text = "→ Validé sous environnements conteneurisés Docker."
        p_d2.font.name = 'Inter'
        p_d2.font.size = Pt(11)
        p_d2.font.color.rgb = DARK_GRAY

    # ============================================================
    # SLIDE 7 : Simulations d'Attaques & Évaluation (10+ lines)
    # ============================================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_background(slide7)
    add_header(slide7, "6. Simulations d'Attaques & Résilience")
    
    # Left Card: Attack Scenario
    draw_card(slide7, Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3), border_color=RED_ACCENT)
    tx_l = slide7.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_l = tx_l.text_frame
    set_text_box_formatting(tf_l)
    p_l = tf_l.paragraphs[0]
    p_l.text = "Scénario : Nœud 3 SilentDrop"
    p_l.font.name = 'Inter'
    p_l.font.size = Pt(18)
    p_l.font.bold = True
    p_l.font.color.rgb = RED_ACCENT
    p_l.space_after = Pt(14)
    
    pts_l = [
        "Déploiement de 4 nœuds isolés via Docker Compose.",
        "Nœud 3 instancié avec le profil SilentNode.",
        "Comportement injecté : rejet complet des communications.",
        "Zéro message de vote prepare ou commit émis.",
        "Simulation parfaite d'un crash réseau complet."
    ]
    for pt in pts_l:
        p_pt = tf_l.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = SILVER
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # Right Card: Real Metrics
    draw_card(slide7, Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3), border_color=GREEN_ACCENT)
    tx_r = slide7.shapes.add_textbox(Inches(6.98), Inches(1.6), Inches(5.6), Inches(5.3))
    tf_r = tx_r.text_frame
    set_text_box_formatting(tf_r)
    p_r = tf_r.paragraphs[0]
    p_r.text = "Résultats de l'Évaluation"
    p_r.font.name = 'Inter'
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = GREEN_ACCENT
    p_r.space_after = Pt(14)
    
    pts_r = [
        "Quorum prepare atteint avec 2/2 prepares.",
        "Quorum commit atteint avec 3/3 commits.",
        "145 rounds successifs complétés sans incident.",
        "Absence complète de forks ou de désynchronisation.",
        "Latence moyenne ultra-faible : 10 ms à 25 ms."
    ]
    for pt in pts_r:
        p_pt = tf_r.add_paragraph()
        p_pt.text = f"• {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = WHITE
        p_pt.space_after = Pt(10)
        p_pt.line_spacing = 1.15

    # ============================================================
    # SLIDE 8 : Le Dashboard Interactif (Futuristic Layout - 10+ lines)
    # ============================================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_background(slide8)
    add_header(slide8, "7. Observabilité en Temps Réel : Dashboard")
    
    draw_card(slide8, Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.3), border_color=CYAN)
    
    txBox = slide8.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.3))
    tf = txBox.text_frame
    set_text_box_formatting(tf)
    p = tf.paragraphs[0]
    p.text = "Architecture d'Observabilité Distribuée"
    p.font.name = 'Inter'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(14)
    
    dashs = [
        "Découplage complet : Application frontend HTML5/CSS3/Javascript statique autonome.",
        "Hébergement : Servie par un conteneur Web Nginx à hautes performances.",
        "Télémétrie : Polling périodique des nœuds via des requêtes HTTP asynchrones légères.",
        "Topologie dynamique : Graphe vectoriel interactif tracé en direct sur Canvas HTML5.",
        "Suivi des rôles : Représentation visuelle animée des nœuds (Leader, Replica, Offline).",
        "Performances en direct : Courbes de latence, débit instantané et volumes de consensus.",
        "Détection d'intrusion : Logs instantanés des alertes en cas de comportement byzantin suspecté."
    ]
    for ds in dashs:
        p_ds = tf.add_paragraph()
        p_ds.text = f"✔ {ds}"
        p_ds.font.name = 'Inter'
        p_ds.font.size = Pt(13)
        p_ds.font.color.rgb = WHITE
        p_ds.space_after = Pt(8)
        p_ds.line_spacing = 1.15

    # ============================================================
    # SLIDE 9 : Conclusion Générale (EMSI Premium Card - 10+ lines)
    # ============================================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_background(slide9)
    add_header(slide9, "Conclusion Générale")
    
    # Elegant full width conclusion card
    draw_card(slide9, Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.3), border_color=PURPLE, border_width=2)
    
    txBox = slide9.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.3))
    tf = txBox.text_frame
    set_text_box_formatting(tf)
    p = tf.paragraphs[0]
    p.text = "Bilan et perspectives de ce projet de PFA à l'EMSI :"
    p.font.name = 'Inter'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.space_after = Pt(14)
    
    conc_pts = [
        "Atteinte complète de tous les objectifs de résilience PBFT fixés.",
        "Mise en pratique approfondie des concepts de multi-threading et de sockets TCP robustes.",
        "Approfondissement des notions cryptographiques modernes NIST et Shamir.",
        "Preuve empirique de la survie du registre malgré la compromission active d'un pair.",
        "Perspectives 1 : Activation de TLS mutuel (mTLS) pour chiffrer l'ensemble des sockets.",
        "Perspectives 2 : Transition vers des signatures de seuil TSS pour optimiser la charge utile.",
        "Perspectives 3 : Simulation de menaces d'équivocation avancée et partitions réseau complexes.",
        "Clôture : Un dossier de PFA académique de haut niveau, prêt pour le jury d'évaluation."
    ]
    for pt in conc_pts:
        p_pt = tf.add_paragraph()
        p_pt.text = f"→ {pt}"
        p_pt.font.name = 'Inter'
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = WHITE
        p_pt.space_after = Pt(8)
        p_pt.line_spacing = 1.15

    # Save the presentation
    filename = "presentation_consensus_byzantin.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == "__main__":
    create_presentation()
